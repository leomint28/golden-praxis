"""
One-shot asset converter for Golden Praxis website.

Run from the project root:
    python convert_assets.py

Requirements (install once):
    pip install pymupdf python-pptx Pillow
"""

import os
import sys

IMAGES = os.path.join(os.path.dirname(os.path.abspath(__file__)), "images")

# ── 1. PDF logo → logo-original.png ──────────────────────────────────────────
def convert_logo_pdf():
    try:
        import fitz  # pymupdf
    except ImportError:
        print("[!] pymupdf not found.  Run:  pip install pymupdf")
        return

    pdf_path = os.path.join(IMAGES, "Golden Praxis Logo.pdf")
    if not os.path.exists(pdf_path):
        print(f"[!] Logo PDF not found at: {pdf_path}")
        return

    doc = fitz.open(pdf_path)
    page = doc[0]

    # 4× scale = ~288 DPI; alpha channel for transparent background
    mat = fitz.Matrix(4, 4)
    pix = page.get_pixmap(matrix=mat, colorspace=fitz.csRGB, alpha=True)
    out_path = os.path.join(IMAGES, "logo-original.png")
    pix.save(out_path)
    print(f"[✓] Logo saved:  {out_path}")
    doc.close()


# ── 2. PPT brand slides → brand-01.png, brand-02.png … ──────────────────────
def extract_ppt_brands():
    # Find any .pptx file in images/
    pptx_files = [
        f for f in os.listdir(IMAGES)
        if f.lower().endswith(".pptx") or f.lower().endswith(".ppt")
    ]
    if not pptx_files:
        print("[!] No .pptx file found in the images folder.")
        print("    Place the brands PowerPoint file inside:  images/")
        return

    pptx_path = os.path.join(IMAGES, pptx_files[0])
    print(f"[i] Processing: {pptx_path}")

    # ── Method A: extract embedded picture shapes ─────────────────────────
    try:
        from pptx import Presentation
        from PIL import Image
        import io
    except ImportError:
        print("[!] python-pptx or Pillow not found.  Run:  pip install python-pptx Pillow")
        return

    prs = Presentation(pptx_path)
    saved = 0
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:           # MSO_SHAPE_TYPE.PICTURE = 13
                blob = shape.image.blob
                img = Image.open(io.BytesIO(blob)).convert("RGBA")
                # Clean name from shape title, fall back to index
                raw_name = (shape.name or f"slide{slide_idx+1}").strip()
                safe_name = "".join(
                    c if c.isalnum() or c in "-_" else "-"
                    for c in raw_name.lower()
                ).strip("-")
                out_path = os.path.join(IMAGES, f"brand-{safe_name}.png")
                img.save(out_path, "PNG")
                print(f"[✓] Brand image saved: {out_path}")
                saved += 1

    if saved == 0:
        print("[i] No embedded picture shapes found — trying slide-thumbnail export …")
        _export_slide_thumbnails(prs, pptx_path)
    else:
        print(f"[✓] {saved} brand image(s) extracted.")


def _export_slide_thumbnails(prs, pptx_path):
    """
    Fallback: render each slide as a PNG via LibreOffice (if installed).
    This handles slides where logos are drawn shapes rather than embedded images.
    """
    import subprocess, glob

    lo_candidates = [
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        "soffice",          # Linux / macOS PATH
        "libreoffice",
    ]
    soffice = None
    for c in lo_candidates:
        if os.path.exists(c) if os.sep in c else True:
            soffice = c
            break

    try:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "png",
             "--outdir", IMAGES, pptx_path],
            check=True, capture_output=True
        )
        # Rename output files to brand-NN.png
        base = os.path.splitext(os.path.basename(pptx_path))[0]
        for i, f in enumerate(sorted(glob.glob(os.path.join(IMAGES, f"{base}*.png")))):
            new_name = os.path.join(IMAGES, f"brand-slide-{i+1:02d}.png")
            os.rename(f, new_name)
            print(f"[✓] Slide thumbnail: {new_name}")
    except Exception as e:
        print(f"[!] LibreOffice export failed: {e}")
        print(
            "\n  If LibreOffice is not installed, export the PPT slides manually:\n"
            "  PowerPoint → File → Export → Change File Type → PNG\n"
            "  Save each slide into the images/ folder as brand-01.png, brand-02.png …\n"
        )


# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("=== Golden Praxis asset converter ===\n")
    convert_logo_pdf()
    print()
    extract_ppt_brands()
    print("\nDone.  Refresh the website to see updated logos.")
